import os
import sys
import PythonMagick
import shutil
from PIL import Image
import PIL.ImageOps
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfFileWriter, PdfFileReader
from logger import Logger
import argparse


class PdfToPpt(object):
    def __init__(self, pdf_file=None, ppt_file=None, crop_regions=None, greyscale=False, blackAndWhite=False, threshold=150, invert=False):
        self.filename = os.path.basename(pdf_file).split('.')[0]
        self.output_dir = self.filename + '_out_dir'

        if os.path.exists(self.output_dir):
            input_var = raw_input(self.output_dir + " exists! Do you want to clear it to continue? (Y/N): ")
            if input_var.lower() != 'y':
                raise Exception('Output directory not available!')

            if not os.path.isdir(self.output_dir):
                os.remove(self.output_dir)
            else:
                shutil.rmtree(self.output_dir)
        
        os.mkdir(self.output_dir)

        self.crop_regions = crop_regions
        self.greyscale = greyscale
        self.blackAndWhite = blackAndWhite
        self.threshold = threshold
        self.invert = invert
        self.pdf_file = pdf_file
        self.ppt_file = os.path.join(self.output_dir, self.filename + '.pptx')
        self.total_pages = 1
        self.image_files = []
        self.log = Logger.defaults('PdfToPptx', os.path.join(self.output_dir, 'out.log'))
        self.log.debug('%s \n %s' % (self.pdf_file, self.ppt_file))

    def check_file_exist(self, file_path):
        self.log.info('Checking file - %s ' % file_path)
        if os.path.isfile(file_path):
            return True
        else:
            return False

    def processing_image(self, image):
        if self.greyscale or self.blackAndWhite:
            image = image.convert('L')

            if self.blackAndWhite:
                pixdata = image.load()
                # Clean the background noise, if color != white, then set to black.
                for x in xrange(image.size[0]):
                    for y in xrange(image.size[1]):
                        pixdata[x, y] = 255 if pixdata[x, y] > self.threshold else 0

        return PIL.ImageOps.invert(image) if self.invert else image

    def pdf_to_image(self, pdf_file):
        if not self.check_file_exist(pdf_file):
            self.log.debug('Requested file not found in %s ' % pdf_file)
            return False
        image_file = pdf_file.replace('.pdf', '.jpg')
        try:
            pdf_to_img = PythonMagick.Image()
            pdf_to_img.density('200')
            pdf_to_img.read(pdf_file)
            pdf_to_img.write(image_file)

            img = Image.open(image_file)
            if self.crop_regions:
                for index, region in enumerate(self.crop_regions):
                    region_image = img.crop(region)
                    region_image = self.processing_image(region_image)
                    region_image_file = image_file.replace('.jpg', '_{}.jpg'.format(index))
                    region_image.save(region_image_file)
                    self.image_files.append(region_image_file)
            else:
                img = self.processing_image(img)
                img.save(image_file)
                self.image_files.append(image_file)

            self.log.info('Image convert passed - %s ' % image_file)
            return True
        except Exception:
            self.log.debug('Image convert failed - %s ' % image_file)
            self.log.error('', exc_info=True)
            return False

    def pdf_splitter(self):
        self.log.info('Called pdf_splitter')
        input_pdf = PdfFileReader(file(self.pdf_file, 'rb'), strict=False)
        self.total_pages = input_pdf.numPages

        for page_number in range(self.total_pages):
            output = PdfFileWriter()
            output.addPage(input_pdf.getPage(page_number))
            # new filename
            new_pdf = os.path.join(self.output_dir, self.filename + '_%s%s' % (str(page_number+1), '.pdf'))
            file_stream = file(new_pdf, 'wb')
            output.write(file_stream)
            file_stream.close()

            # calling pdf to image conversion
            self.pdf_to_image(new_pdf)

    def create_ppt(self):
        self.log.info('Called create_ppt')
        prs = Presentation()
        try:
            for slide_number, img_path in enumerate(self.image_files):
                self.log.debug('%s' % img_path)
                new_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(new_slide)
                subtitle = slide.placeholders[1]
                title = slide.shapes.title
                title.text = "Image %s " % str(slide_number+1)
                left = top = Inches(0)
                height = Inches(7.5)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                prs.save(self.ppt_file)
        except IOError:
            self.log.error('error creating ppt', exc_info=True)

    def execute(self):
        self.log.info('Calling the main execution for ppt conversion')
        self.pdf_splitter()
        self.create_ppt()
        self.log.info('Done ppt conversion')

def region(s):
    try:
        left, top, right, bottom = map(int, s.split(','))
        return left, top, right, bottom
    except:
        raise argparse.ArgumentTypeError("Coordinates must be left,top,right,bottom")

if __name__ == '__main__':
     #run_time = sys.argv[1:]

    parser = argparse.ArgumentParser(description='Convert PDF to pptx')
    parser.add_argument('pdf_file',
                        help='path to the pdf file')
    parser.add_argument('-g', '--greyscale', action='store_true',
                        help='generate greyscale pptx')
    parser.add_argument('-b', '--blackwhite', action='store_true',
                        help='generate black and white pptx')
    parser.add_argument('-t', '--threshold', type=int, choices=range(0,256), default=150,
                        help='the threshold value [0,255) for converting black and white')
    parser.add_argument('-i', '--invert', action='store_true',
                        help='invert generated pptx color')
    parser.add_argument('-c', '--crop', type=region, nargs='+',
                        help='values specifing region(s) in pixel to crop')
    args = parser.parse_args()
    print args
    PdfToPpt(pdf_file=args.pdf_file, greyscale=args.greyscale,blackAndWhite=args.blackwhite,
            threshold=args.threshold, invert=args.invert, crop_regions=args.crop).execute()

